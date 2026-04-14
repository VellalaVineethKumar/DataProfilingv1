streanerate a business-oriented PowerPoint presentation for Master Data Profiler."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path

# ── Brand palette ──────────────────────────────────────────────────────────
# Primary: #B60003 | Supporting: #221E1F
DARK_BG  = RGBColor(0x22, 0x1E, 0x1F)      # supporting dark (slide bg)
ACCENT   = RGBColor(0xB6, 0x00, 0x03)       # primary red (headlines, bars, CTAs)
ACCENT2  = RGBColor(0xED, 0x55, 0x56)       # primary light 1 (secondary highlights)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xBE, 0xBE, 0xBD)       # body text on dark backgrounds
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)       # warning / stat callout
CARD_BG  = RGBColor(0x32, 0x2E, 0x2F)       # card panels (lighter than bg)
GREEN    = RGBColor(0x10, 0xB9, 0x81)       # success / positive indicators
PINK_LT  = RGBColor(0xF1, 0x8E, 0x8E)       # primary light 3 (subtle accent)
MUTED    = RGBColor(0x76, 0x75, 0x75)       # supporting mid 2 (captions)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def _solid_bg(slide, color: RGBColor):
    """Fill the slide background with a solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape(slide, left, top, width, height, fill_color: RGBColor, corner_radius=Inches(0.15)):
    """Add a rounded-rectangle card shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    if hasattr(shape, "adjustments") and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.05
    return shape


def _text_box(slide, left, top, width, height, text, font_size=18, bold=False,
              color=WHITE, alignment=PP_ALIGN.LEFT, font_name="Segoe UI"):
    """Add a text box with common formatting."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _bullet_frame(slide, left, top, width, height, items, font_size=16,
                  color=LIGHT, bullet_color=ACCENT, font_name="Segoe UI"):
    """Add a text frame with bullet points."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_before = Pt(6)
        p.space_after = Pt(4)
        p.level = 0
    return txBox


def _section_divider(prs, title_text, subtitle_text=""):
    """Create a bold section-divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _solid_bg(slide, DARK_BG)

    # accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(3.1), SLIDE_W, Inches(1.3))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()

    _text_box(slide, Inches(0.8), Inches(3.2), Inches(11), Inches(0.7),
              title_text, font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    if subtitle_text:
        _text_box(slide, Inches(0.8), Inches(3.85), Inches(11), Inches(0.5),
                  subtitle_text, font_size=18, color=WHITE, alignment=PP_ALIGN.CENTER)
    return slide


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ═══════════════════════════════════════════════════════════════════════════

def slide_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)

    # decorative accent stripe
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()

    _text_box(slide, Inches(1.2), Inches(1.8), Inches(10), Inches(1),
              "MASTER DATA PROFILER", font_size=48, bold=True, color=WHITE)
    _text_box(slide, Inches(1.2), Inches(2.9), Inches(10), Inches(0.8),
              "Automated Data Quality & Profiling Platform", font_size=24, color=ACCENT)
    _text_box(slide, Inches(1.2), Inches(4.2), Inches(10), Inches(0.6),
              "Turn raw data into trusted, actionable insights — in minutes, not months.",
              font_size=16, color=LIGHT)
    _text_box(slide, Inches(1.2), Inches(6.2), Inches(10), Inches(0.4),
              "Business Overview  |  April 2026", font_size=14, color=LIGHT)


def slide_agenda(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Agenda", font_size=36, bold=True, color=WHITE)

    items = [
        ("01", "The Problem — Why Data Quality Matters"),
        ("02", "Product Overview & Key Capabilities"),
        ("03", "Live Dashboard & Profiling Features"),
        ("04", "Data Quality Engine & AI Recommendations"),
        ("05", "Enterprise Connectivity (15+ Databases)"),
        ("06", "Reporting, Export & Collaboration"),
        ("07", "Architecture & Deployment"),
        ("08", "Security & Governance"),
        ("09", "Roadmap & Future Vision"),
        ("10", "Next Steps"),
    ]
    for i, (num, label) in enumerate(items):
        y = Inches(1.5) + Inches(i * 0.55)
        _text_box(slide, Inches(1.0), y, Inches(0.6), Inches(0.45),
                  num, font_size=20, bold=True, color=ACCENT)
        _text_box(slide, Inches(1.7), y, Inches(9), Inches(0.45),
                  label, font_size=18, color=LIGHT)


def slide_problem(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "The Problem", font_size=36, bold=True, color=WHITE)
    _text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
              "Bad data costs the average enterprise $12.9 M per year  (Gartner)",
              font_size=16, color=AMBER, bold=True)

    pain_points = [
        "Manual profiling is slow — analysts spend 40-60% of their time cleaning data",
        "Inconsistent quality checks across teams lead to conflicting reports",
        "No single pane of glass for data health across multiple sources",
        "Regulatory non-compliance risk when data lineage is unclear",
        "Data drift goes undetected until downstream models break",
    ]
    _bullet_frame(slide, Inches(0.8), Inches(2.0), Inches(11), Inches(3.5),
                  [f"•  {p}" for p in pain_points], font_size=17, color=LIGHT)

    _add_shape(slide, Inches(0.8), Inches(5.7), Inches(11.5), Inches(1.2), CARD_BG)
    _text_box(slide, Inches(1.2), Inches(5.85), Inches(10.8), Inches(0.9),
              "\"We need a platform that automates data profiling, enforces quality rules, "
              "and provides real-time visibility — without requiring a data engineering team.\"",
              font_size=15, color=LIGHT, alignment=PP_ALIGN.CENTER)


def slide_overview(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Product Overview", font_size=36, bold=True, color=WHITE)
    _text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
              "Master Data Profiler is a self-service data quality & profiling platform built for speed and simplicity.",
              font_size=16, color=LIGHT)

    cards = [
        ("Instant Profiling", "Upload CSV, Excel, Parquet, JSON, Feather or connect\nto 15+ databases. Profile in seconds."),
        ("Quality Rules Engine", "50+ built-in checks. Import/export rules.\nAI-powered recommendations."),
        ("Live Dashboard", "KPI cards, quality gauge, null analysis,\ndata-type distribution — all at a glance."),
        ("Enterprise Ready", "Role-based auth, audit trail, Docker\ndeployment, REST API, PDF/HTML reports."),
    ]
    for i, (title, desc) in enumerate(cards):
        x = Inches(0.6) + Inches(i * 3.1)
        _add_shape(slide, x, Inches(2.2), Inches(2.9), Inches(3.8), CARD_BG)
        _text_box(slide, x + Inches(0.2), Inches(2.4), Inches(2.5), Inches(0.5),
                  title, font_size=18, bold=True, color=ACCENT)
        _text_box(slide, x + Inches(0.2), Inches(3.1), Inches(2.5), Inches(2.8),
                  desc, font_size=14, color=LIGHT)


def slide_profiling_features(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Profiling Capabilities", font_size=36, bold=True, color=WHITE)

    left_items = [
        "Column-level statistics (mean, median, mode, std, min, max, quartiles)",
        "Null / missing-value analysis with visual heatmaps",
        "Unique-value & cardinality analysis",
        "Data-type detection and distribution charts",
        "Cross-column correlation matrix (Pearson)",
        "Duplicate row detection with configurable thresholds",
    ]
    right_items = [
        "Pattern analysis for string columns (email, phone, dates)",
        "Outlier detection (IQR & z-score methods)",
        "Statistical distribution identification",
        "Large-file support with streaming & chunked loading",
        "Multi-file comparison across datasets",
        "Data drift detection against saved baselines",
    ]

    _add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), CARD_BG)
    _bullet_frame(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(5.0),
                  [f"✓  {p}" for p in left_items], font_size=14, color=LIGHT)

    _add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3), CARD_BG)
    _bullet_frame(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.0),
                  [f"✓  {p}" for p in right_items], font_size=14, color=LIGHT)


def slide_dashboard(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Interactive Dashboard", font_size=36, bold=True, color=WHITE)
    _text_box(slide, Inches(0.8), Inches(1.2), Inches(10), Inches(0.5),
              "A single pane of glass for data health — designed for analysts and leaders alike.",
              font_size=16, color=LIGHT)

    kpis = [
        ("Rows & Columns", "Instant shape overview"),
        ("Quality Score", "0-100 gauge updated in real time"),
        ("Top Issues", "Missing values, outliers, type mismatches"),
        ("Null Heatmap", "Column-level null percentages"),
        ("Correlation View", "Highlight highly-correlated pairs"),
        ("Recent Operations", "Audit trail of actions taken"),
    ]
    for i, (title, desc) in enumerate(kpis):
        col = i % 3
        row = i // 3
        x = Inches(0.6) + Inches(col * 4.1)
        y = Inches(2.2) + Inches(row * 2.4)
        _add_shape(slide, x, y, Inches(3.8), Inches(2.0), CARD_BG)
        _text_box(slide, x + Inches(0.25), y + Inches(0.25), Inches(3.3), Inches(0.45),
                  title, font_size=17, bold=True, color=ACCENT2)
        _text_box(slide, x + Inches(0.25), y + Inches(0.8), Inches(3.3), Inches(1.0),
                  desc, font_size=14, color=LIGHT)


def slide_quality_engine(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Data Quality Engine", font_size=36, bold=True, color=WHITE)

    _add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), CARD_BG)
    _text_box(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.5),
              "Built-in Rule Types", font_size=20, bold=True, color=ACCENT)
    rules = [
        "Not-null / completeness checks",
        "Uniqueness constraints",
        "Range & boundary validation",
        "Regex pattern matching",
        "Allowed-value (enum) lists",
        "Cross-column consistency",
        "Custom SQL expressions",
    ]
    _bullet_frame(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.2),
                  [f"•  {r}" for r in rules], font_size=14, color=LIGHT)

    _add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3), CARD_BG)
    _text_box(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5),
              "AI-Powered Recommendations", font_size=20, bold=True, color=ACCENT)
    ai_items = [
        "Azure OpenAI integration for rule suggestions",
        "Heuristic fallback when LLM is unavailable",
        "Auto-detect email, phone, date patterns",
        "Smart threshold calibration from sample data",
        "One-click accept & apply suggested rules",
        "Rule library — save, share & reuse across teams",
        "Import / export rules as JSON",
    ]
    _bullet_frame(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.2),
                  [f"•  {r}" for r in ai_items], font_size=14, color=LIGHT)


def slide_databases(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Enterprise Connectivity", font_size=36, bold=True, color=WHITE)
    _text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
              "Connect to 15+ databases via SQLAlchemy — one unified interface, zero vendor lock-in.",
              font_size=16, color=LIGHT)

    dbs = [
        "PostgreSQL", "MySQL", "MariaDB", "SQL Server", "Oracle",
        "SQLite", "IBM Db2", "Snowflake", "BigQuery", "Redshift",
        "CockroachDB", "DuckDB", "Teradata", "SAP HANA", "Apache Hive",
    ]
    for i, db in enumerate(dbs):
        col = i % 5
        row = i // 5
        x = Inches(0.5) + Inches(col * 2.5)
        y = Inches(2.3) + Inches(row * 1.5)
        _add_shape(slide, x, y, Inches(2.2), Inches(1.1), CARD_BG)
        _text_box(slide, x + Inches(0.15), y + Inches(0.25), Inches(1.9), Inches(0.6),
                  db, font_size=15, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)

    _text_box(slide, Inches(0.8), Inches(6.2), Inches(11), Inches(0.5),
              "Supports raw SQLAlchemy URLs  •  Browse schemas, tables & views  •  Run custom SQL queries",
              font_size=14, color=LIGHT, alignment=PP_ALIGN.CENTER)


def slide_file_formats(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "File Format Support & Multi-File Analysis", font_size=36, bold=True, color=WHITE)

    _add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), CARD_BG)
    _text_box(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.5),
              "Supported Formats", font_size=20, bold=True, color=ACCENT)
    formats = [
        "CSV  — with delimiter & encoding auto-detection",
        "Excel (.xlsx, .xls)  — multi-sheet support",
        "JSON  — nested & flat structures",
        "JSONL  — streaming line-delimited JSON",
        "Parquet  — columnar analytics format",
        "Feather  — fast in-memory interchange",
    ]
    _bullet_frame(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.0),
                  [f"•  {f}" for f in formats], font_size=14, color=LIGHT)

    _add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3), CARD_BG)
    _text_box(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5),
              "Multi-File Comparison", font_size=20, bold=True, color=ACCENT)
    multi_items = [
        "Upload multiple files side-by-side",
        "Schema alignment & diff view",
        "Row-count & null-percentage comparison",
        "Common-column bar charts",
        "Identify discrepancies across data sources",
        "Ideal for migration validation & ETL QA",
    ]
    _bullet_frame(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.0),
                  [f"•  {m}" for m in multi_items], font_size=14, color=LIGHT)


def slide_reporting(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Reporting & Export", font_size=36, bold=True, color=WHITE)

    items = [
        ("PDF Reports", "Comprehensive profiling reports with charts, scores, and issue tables — ready for management review."),
        ("HTML Reports", "Shareable single-file HTML reports that open in any browser."),
        ("Data Export", "Download cleaned data as CSV, Excel, JSON, Parquet, or Feather."),
        ("Rule Export", "Export / import quality rule configurations as JSON for team sharing."),
        ("REST API", "FastAPI endpoint for programmatic profiling — integrate with CI/CD pipelines."),
    ]
    for i, (title, desc) in enumerate(items):
        y = Inches(1.6) + Inches(i * 1.1)
        _add_shape(slide, Inches(0.5), y, Inches(12.2), Inches(0.9), CARD_BG)
        _text_box(slide, Inches(0.7), y + Inches(0.1), Inches(2.5), Inches(0.4),
                  title, font_size=17, bold=True, color=ACCENT)
        _text_box(slide, Inches(3.2), y + Inches(0.15), Inches(9.2), Inches(0.65),
                  desc, font_size=14, color=LIGHT)


def slide_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Architecture & Deployment", font_size=36, bold=True, color=WHITE)

    layers = [
        ("Presentation Layer", "Streamlit interactive UI  •  Responsive layout  •  Plotly charts",  ACCENT),
        ("API Layer",           "FastAPI REST endpoints  •  Pydantic models  •  Health checks",      ACCENT2),
        ("Core Engine",         "Pandas profiler  •  Quality rules  •  Drift detector  •  AI recs",  AMBER),
        ("Persistence",         "SQLite (rules, audit, projects)  •  File system  •  Session state", PINK_LT),
        ("Connectivity",        "SQLAlchemy  •  15+ DB drivers  •  File parsers (6 formats)",        ACCENT),
    ]
    for i, (layer, desc, color) in enumerate(layers):
        y = Inches(1.5) + Inches(i * 1.1)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), y, Inches(0.12), Inches(0.9))
        bar.fill.solid(); bar.fill.fore_color.rgb = color; bar.line.fill.background()
        _add_shape(slide, Inches(0.8), y, Inches(11.8), Inches(0.9), CARD_BG)
        _text_box(slide, Inches(1.0), y + Inches(0.08), Inches(3), Inches(0.4),
                  layer, font_size=16, bold=True, color=color)
        _text_box(slide, Inches(1.0), y + Inches(0.45), Inches(11.2), Inches(0.4),
                  desc, font_size=13, color=LIGHT)

    _text_box(slide, Inches(0.8), Inches(6.7), Inches(11), Inches(0.4),
              "Docker & Docker Compose ready  •  Single-command deployment  •  Health-check built in",
              font_size=14, color=GREEN, alignment=PP_ALIGN.CENTER)


def slide_security(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Security & Governance", font_size=36, bold=True, color=WHITE)

    _add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.3), CARD_BG)
    _text_box(slide, Inches(0.7), Inches(1.7), Inches(5.5), Inches(0.5),
              "Access Control", font_size=20, bold=True, color=ACCENT)
    sec_items = [
        "User authentication with salted SHA-256 password hashing",
        "Role-based access (admin / analyst)",
        "Session management with automatic timeout",
        "Secrets management via Streamlit secrets.toml",
        "No plaintext passwords stored — automatic migration",
    ]
    _bullet_frame(slide, Inches(0.7), Inches(2.3), Inches(5.5), Inches(4.0),
                  [f"•  {s}" for s in sec_items], font_size=14, color=LIGHT)

    _add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.9), Inches(5.3), CARD_BG)
    _text_box(slide, Inches(7.0), Inches(1.7), Inches(5.5), Inches(0.5),
              "Audit & Compliance", font_size=20, bold=True, color=ACCENT)
    audit_items = [
        "Full audit trail — every action logged with timestamp & user",
        "Persistent SQLite-backed audit log",
        "Audit log visible in sidebar for real-time review",
        "Project workspaces for data isolation",
        "Export audit records for compliance reporting",
    ]
    _bullet_frame(slide, Inches(7.0), Inches(2.3), Inches(5.5), Inches(4.0),
                  [f"•  {a}" for a in audit_items], font_size=14, color=LIGHT)


def slide_drift(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Data Drift Detection", font_size=36, bold=True, color=WHITE)
    _text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.5),
              "Catch silent data degradation before it impacts downstream models and reports.",
              font_size=16, color=LIGHT)

    steps = [
        ("1. Save Baseline", "Capture today's profile as the reference point for any dataset."),
        ("2. Re-Profile Later", "Load new data (daily feed, monthly refresh, etc.) and re-run profiling."),
        ("3. Detect Drift", "Automatic comparison — null-rate shifts, mean changes, cardinality drift."),
        ("4. Alert & Act", "Visual drift report highlights columns that moved beyond your threshold."),
    ]
    for i, (step, desc) in enumerate(steps):
        x = Inches(0.4) + Inches(i * 3.2)
        _add_shape(slide, x, Inches(2.3), Inches(2.9), Inches(3.5), CARD_BG)
        _text_box(slide, x + Inches(0.2), Inches(2.5), Inches(2.5), Inches(0.5),
                  step, font_size=17, bold=True, color=ACCENT2)
        _text_box(slide, x + Inches(0.2), Inches(3.2), Inches(2.5), Inches(2.4),
                  desc, font_size=14, color=LIGHT)


def slide_competitive(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Why Master Data Profiler?", font_size=36, bold=True, color=WHITE)

    rows = [
        ("Capability",        "Master Data Profiler",   "Typical Tools"),
        ("Setup time",        "< 5 minutes (Docker)",   "Days to weeks"),
        ("Database support",  "15+ via SQLAlchemy",     "3-5 connectors"),
        ("AI rule suggestions","Built-in (Azure OpenAI)","Add-on / none"),
        ("Data drift",        "Included",               "Separate product"),
        ("Audit trail",       "Built-in SQLite log",    "Manual / external"),
        ("PDF / HTML reports","One-click export",       "Custom scripting"),
        ("Self-hosted",       "Yes (Docker / bare-metal)","Cloud-only lock-in"),
        ("Pricing model",     "Open / internal tool",   "$$$$ per seat"),
    ]

    col_widths = [Inches(2.8), Inches(4.2), Inches(4.2)]
    x_starts = [Inches(0.6), Inches(3.4), Inches(7.6)]
    for ri, (c1, c2, c3) in enumerate(rows):
        y = Inches(1.6) + Inches(ri * 0.6)
        bg = CARD_BG if ri > 0 else ACCENT
        text_c = LIGHT if ri > 0 else WHITE
        for ci, val in enumerate([c1, c2, c3]):
            clr = GREEN if (ci == 1 and ri > 0) else text_c
            _add_shape(slide, x_starts[ci], y, col_widths[ci] - Inches(0.1), Inches(0.5), bg)
            _text_box(slide, x_starts[ci] + Inches(0.15), y + Inches(0.05),
                      col_widths[ci] - Inches(0.3), Inches(0.4),
                      val, font_size=13, bold=(ri == 0), color=clr if ri > 0 else WHITE,
                      alignment=PP_ALIGN.CENTER if ci > 0 else PP_ALIGN.LEFT)


def slide_roadmap(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Roadmap & Vision", font_size=36, bold=True, color=WHITE)

    quarters = [
        ("Q2 2026 — Now", [
            "15+ database connectors",
            "AI rule recommendations",
            "Data drift detection",
            "PDF / HTML reporting",
            "Docker deployment",
        ]),
        ("Q3 2026", [
            "Scheduled profiling jobs",
            "Email / Slack alerting",
            "Column-level lineage",
            "Enhanced regex library",
            "Team collaboration features",
        ]),
        ("Q4 2026", [
            "Automated remediation",
            "ML-powered anomaly detection",
            "Custom dashboards builder",
            "SSO / LDAP integration",
            "Cloud-native deployment (K8s)",
        ]),
        ("2027+", [
            "Federated profiling (multi-node)",
            "Real-time streaming profiling",
            "Data catalog integration",
            "Marketplace for community rules",
            "Mobile-friendly dashboard",
        ]),
    ]
    for i, (title, items) in enumerate(quarters):
        x = Inches(0.3) + Inches(i * 3.2)
        _add_shape(slide, x, Inches(1.6), Inches(3.0), Inches(5.2), CARD_BG)
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(1.6), Inches(3.0), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT if i == 0 else ACCENT2
        bar.line.fill.background()
        _text_box(slide, x + Inches(0.2), Inches(1.8), Inches(2.6), Inches(0.5),
                  title, font_size=16, bold=True, color=ACCENT if i == 0 else PINK_LT)
        _bullet_frame(slide, x + Inches(0.2), Inches(2.4), Inches(2.6), Inches(4.2),
                      [f"•  {it}" for it in items], font_size=12, color=LIGHT)


def slide_next_steps(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)
    _text_box(slide, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
              "Next Steps", font_size=36, bold=True, color=WHITE)

    steps = [
        ("1", "Live Demo", "Schedule a 30-minute walkthrough with your data team to see Master Data Profiler in action on your own datasets."),
        ("2", "Pilot Program", "Deploy on a single use-case (e.g., monthly financial close data) to quantify time savings and quality improvement."),
        ("3", "Rollout & Scale", "Expand to additional teams and data sources. Leverage Docker for enterprise-wide deployment."),
    ]
    for i, (num, title, desc) in enumerate(steps):
        y = Inches(1.7) + Inches(i * 1.7)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(0.8), y + Inches(0.1), Inches(0.7), Inches(0.7))
        circle.fill.solid(); circle.fill.fore_color.rgb = ACCENT; circle.line.fill.background()
        _text_box(slide, Inches(0.85), y + Inches(0.18), Inches(0.6), Inches(0.5),
                  num, font_size=24, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
        _text_box(slide, Inches(1.8), y, Inches(10), Inches(0.5),
                  title, font_size=22, bold=True, color=ACCENT2)

        _text_box(slide, Inches(1.8), y + Inches(0.55), Inches(10), Inches(0.8),
                  desc, font_size=15, color=LIGHT)

    _text_box(slide, Inches(0.8), Inches(6.5), Inches(11), Inches(0.5),
              "Let's turn your data into your competitive advantage.",
              font_size=18, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)


def slide_thank_you(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _solid_bg(slide, DARK_BG)

    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.15), SLIDE_H)
    stripe.fill.solid(); stripe.fill.fore_color.rgb = ACCENT; stripe.line.fill.background()

    _text_box(slide, Inches(1.2), Inches(2.5), Inches(10), Inches(1),
              "Thank You", font_size=52, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _text_box(slide, Inches(1.2), Inches(3.7), Inches(10), Inches(0.6),
              "Master Data Profiler", font_size=28, color=ACCENT, alignment=PP_ALIGN.CENTER)
    _text_box(slide, Inches(1.2), Inches(4.6), Inches(10), Inches(0.6),
              "Questions?  •  Demo?  •  Let's Talk.", font_size=20, color=LIGHT, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════
# MAIN
# ═══════════════════════════════════════════════════════════════════════════

def build_presentation() -> Path:
    """Build the full presentation and return the output path."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)                          # 1
    slide_agenda(prs)                         # 2
    slide_problem(prs)                        # 3
    _section_divider(prs, "PRODUCT DEEP DIVE", "Capabilities that set us apart")  # 4
    slide_overview(prs)                       # 5
    slide_profiling_features(prs)             # 6
    slide_dashboard(prs)                      # 7
    slide_quality_engine(prs)                 # 8
    slide_databases(prs)                      # 9
    slide_file_formats(prs)                   # 10
    slide_drift(prs)                          # 11
    slide_reporting(prs)                      # 12
    _section_divider(prs, "ENTERPRISE GRADE", "Architecture  •  Security  •  Deployment")  # 13
    slide_architecture(prs)                   # 14
    slide_security(prs)                       # 15
    slide_competitive(prs)                    # 16
    slide_roadmap(prs)                        # 17
    slide_next_steps(prs)                     # 18
    slide_thank_you(prs)                      # 19

    out = Path("Master_Data_Profiler_Branded.pptx")
    prs.save(str(out))
    return out


if __name__ == "__main__":
    path = build_presentation()
    print(f"Presentation saved to: {path}")
    print(f"Total slides: 19")
